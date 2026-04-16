from __future__ import annotations

import io

from pptx import Presentation
from pptx.util import Pt


def extract_slide_titles_from_pptx(pptx_source: bytes | io.BytesIO) -> str:
    """
    Extract the title text from each slide's title placeholder.
    Returns a newline-separated string of titles (one per line), suitable
    for pasting directly into the outline text area.
    Falls back to "Slide N" when a slide has no title placeholder.
    """
    buf = io.BytesIO(pptx_source) if isinstance(pptx_source, bytes) else pptx_source
    prs = Presentation(buf)

    titles: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        title: str | None = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.idx == 0:
                title = shape.text_frame.text.strip()
                break
        titles.append(title or f"Slide {i}")

    return "\n".join(titles)


def write_notes_to_pptx(
    pptx_source: bytes | io.BytesIO,
    notes: list[dict],  # [{"number": 1, "title": "...", "notes": "..."}]
) -> bytes:
    """
    Write speaker notes into a copy of the provided .pptx file.
    Matching is index-based: notes[0] → slide 1, notes[1] → slide 2, etc.
    The original file is never modified; the result is returned as bytes.
    """
    buf = io.BytesIO(pptx_source) if isinstance(pptx_source, bytes) else pptx_source
    prs = Presentation(buf)

    notes_by_number: dict[int, str] = {n["number"]: n["notes"] for n in notes}

    for i, slide in enumerate(prs.slides, start=1):
        note_text = notes_by_number.get(i, "")
        if not note_text:
            continue

        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame

        # Remove extra paragraphs, keeping the first (notes placeholder paragraph)
        paras = tf.paragraphs
        for para in list(paras)[1:]:
            p_elem = para._p
            p_elem.getparent().remove(p_elem)

        tf.paragraphs[0].text = note_text
        for run in tf.paragraphs[0].runs:
            run.font.size = Pt(12)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
