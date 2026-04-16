from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.util import Pt

from speakernotes_whisper.pptx_writer import (
    extract_slide_titles_from_pptx,
    write_notes_to_pptx,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_pptx_bytes(*titles: str) -> bytes:
    """Return bytes of an in-memory pptx with one titled slide per title."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content — guaranteed title placeholder (idx=0)
    for title in titles:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_blank_pptx_bytes(n_slides: int = 1) -> bytes:
    """Return bytes of an in-memory pptx with blank slides (no title placeholder)."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank — no placeholders
    for _ in range(n_slides):
        prs.slides.add_slide(blank_layout)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _read_notes(pptx_bytes: bytes, slide_index: int) -> str:
    """Extract the notes text from the given slide (0-indexed) of a pptx bytes object."""
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide = prs.slides[slide_index]
    return slide.notes_slide.notes_text_frame.text.strip()


# ---------------------------------------------------------------------------
# extract_slide_titles_from_pptx
# ---------------------------------------------------------------------------


class TestExtractSlideTitles:
    def test_single_slide_title(self):
        result = extract_slide_titles_from_pptx(_make_pptx_bytes("My Title"))
        assert result == "My Title"

    def test_multiple_slide_titles_newline_separated(self):
        result = extract_slide_titles_from_pptx(_make_pptx_bytes("First", "Second", "Third"))
        assert result == "First\nSecond\nThird"

    def test_accepts_bytes_input(self):
        data = _make_pptx_bytes("Bytes Title")
        assert isinstance(data, bytes)
        result = extract_slide_titles_from_pptx(data)
        assert result == "Bytes Title"

    def test_accepts_bytesio_input(self):
        data = _make_pptx_bytes("BytesIO Title")
        result = extract_slide_titles_from_pptx(io.BytesIO(data))
        assert result == "BytesIO Title"

    def test_blank_slide_falls_back_to_slide_n(self):
        result = extract_slide_titles_from_pptx(_make_blank_pptx_bytes(1))
        assert result == "Slide 1"

    def test_multiple_blank_slides_numbered_correctly(self):
        result = extract_slide_titles_from_pptx(_make_blank_pptx_bytes(3))
        assert result == "Slide 1\nSlide 2\nSlide 3"

    def test_mixed_titled_and_blank_slides(self):
        prs = Presentation()
        titled_layout = prs.slide_layouts[1]
        blank_layout = prs.slide_layouts[6]

        slide1 = prs.slides.add_slide(titled_layout)
        slide1.shapes.title.text = "Real Title"
        prs.slides.add_slide(blank_layout)  # slide 2 has no title

        buf = io.BytesIO()
        prs.save(buf)

        result = extract_slide_titles_from_pptx(buf.getvalue())
        lines = result.split("\n")
        assert lines[0] == "Real Title"
        assert lines[1] == "Slide 2"

    def test_title_text_is_stripped(self):
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "  Padded Title  "
        buf = io.BytesIO()
        prs.save(buf)

        result = extract_slide_titles_from_pptx(buf.getvalue())
        assert result == "Padded Title"


# ---------------------------------------------------------------------------
# write_notes_to_pptx
# ---------------------------------------------------------------------------


class TestWriteNotesToPptx:
    def test_returns_bytes(self):
        pptx = _make_pptx_bytes("Slide 1")
        result = write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "Hello"}])
        assert isinstance(result, bytes)

    def test_notes_written_to_correct_slide(self):
        pptx = _make_pptx_bytes("Slide 1")
        result = write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "My speaker notes"}])
        assert _read_notes(result, 0) == "My speaker notes"

    def test_notes_written_to_multiple_slides(self):
        pptx = _make_pptx_bytes("Slide 1", "Slide 2", "Slide 3")
        notes = [
            {"number": 1, "title": "Slide 1", "notes": "Notes for one"},
            {"number": 2, "title": "Slide 2", "notes": "Notes for two"},
            {"number": 3, "title": "Slide 3", "notes": "Notes for three"},
        ]
        result = write_notes_to_pptx(pptx, notes)
        assert _read_notes(result, 0) == "Notes for one"
        assert _read_notes(result, 1) == "Notes for two"
        assert _read_notes(result, 2) == "Notes for three"

    def test_slide_with_no_matching_note_is_unchanged(self):
        pptx = _make_pptx_bytes("Slide 1", "Slide 2")
        # Only provide a note for slide 1; slide 2 should remain empty
        result = write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "Note"}])
        assert _read_notes(result, 1) == ""

    def test_original_bytes_are_not_modified(self):
        pptx = _make_pptx_bytes("Slide 1")
        original_copy = bytes(pptx)
        write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "Notes"}])
        assert pptx == original_copy

    def test_accepts_bytes_input(self):
        pptx = _make_pptx_bytes("Slide 1")
        assert isinstance(pptx, bytes)
        result = write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "Note"}])
        assert _read_notes(result, 0) == "Note"

    def test_accepts_bytesio_input(self):
        pptx = _make_pptx_bytes("Slide 1")
        result = write_notes_to_pptx(io.BytesIO(pptx), [{"number": 1, "title": "Slide 1", "notes": "Note"}])
        assert _read_notes(result, 0) == "Note"

    def test_empty_notes_list_returns_valid_pptx(self):
        pptx = _make_pptx_bytes("Slide 1")
        result = write_notes_to_pptx(pptx, [])
        # Should still be parseable
        Presentation(io.BytesIO(result))

    def test_note_number_beyond_slide_count_is_ignored(self):
        pptx = _make_pptx_bytes("Slide 1")
        # Number 99 doesn't correspond to any slide — should not raise
        result = write_notes_to_pptx(
            pptx,
            [
                {"number": 1, "title": "Slide 1", "notes": "Good note"},
                {"number": 99, "title": "Ghost", "notes": "No slide"},
            ],
        )
        assert _read_notes(result, 0) == "Good note"

    def test_font_size_set_to_12pt(self):
        pptx = _make_pptx_bytes("Slide 1")
        result = write_notes_to_pptx(pptx, [{"number": 1, "title": "Slide 1", "notes": "Check font"}])

        prs = Presentation(io.BytesIO(result))
        slide = prs.slides[0]
        para = slide.notes_slide.notes_text_frame.paragraphs[0]
        sizes = [run.font.size for run in para.runs if run.font.size is not None]
        assert sizes, "Expected at least one run with an explicit font size"
        assert all(s == Pt(12) for s in sizes)
